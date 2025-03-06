from langchain_google_genai import GoogleGenerativeAIEmbeddings
from langchain_community.vectorstores import Chroma
from langchain_google_genai import GoogleGenerativeAI
from langchain_ollama.llms import OllamaLLM
from langchain_ollama import OllamaEmbeddings
from pptx import Presentation
from langchain_core.documents import Document
from langchain_community.document_loaders import PDFPlumberLoader
from langchain_text_splitters import RecursiveCharacterTextSplitter
from langchain_core.prompts import ChatPromptTemplate
from typing import List
import pandas as pd

class LanguageModel:
    def __init__(self, model: str = "gemini-2.0-flash") -> None:
        """
        Initialize the LanguageModel class.

        Args:
            model (str, optional): The language model to use. Defaults to "gemini-2.0-flash".

        Raises:
            ValueError: If the model is not recognized.

        """
        
        if "gemini" in model:
            self.model = GoogleGenerativeAI(model=model)
            self.embedding = GoogleGenerativeAIEmbeddings(model="models/embedding-001")
        else:
            self.model = OllamaLLM(model=model)
            self.embedding = OllamaEmbeddings(model=model)

        self.document_vector_db = Chroma(
            collection_name="doc_store",
            embedding_function=self.embedding,
            persist_directory="../data/chroma_db"
        )

        self.PROMPT_TEMPLATE = """\
        You are an expert assistant. Answer the query directly using the provided context.
        If unsure, state that you don't know. Always answer in Bahasa Indonesia.
        
        Do not start your response with phrases like "Berdasarkan konteks yang diberikan" or similar introductory statements.
        
        Query: {user_query} 
        Context: {document_context} 
        Answer:
        """

        
    def load_pptx(self, path: str) -> str:
        """
        Load a PowerPoint (.pptx) file and extract all text from it into a single string.

        Args:
            path (str): The path to the PowerPoint file.

        Returns:
            List[Document]: A list of Document objects containing the extracted text in a single string.
        """
        prs = Presentation(path)
        texts = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    texts.append(shape.text)
        combined_text = " ".join(texts)
        return [Document(page_content=combined_text, metadata={ "source" : path })]
    
    def load_pdf(self, path: str) -> List[any]:
        """
        Load a PDF file and extract all text from it into a list of Document objects.

        Args:
            path (str): The path to the PDF file.

        Returns:
            List[Document]: A list of Document objects containing the extracted text.
        """
        doc_loader = PDFPlumberLoader(path)
        return doc_loader.load()
    
    def load_xlsx(self, path: str) -> List[Document]:
        """
        Load an Excel (.xlsx) file and extract all text from it into a single string.

        The Excel file is read into a pandas DataFrame and then converted into a CSV string.
        The resulting string is returned as a list of Document objects.

        Args:
            path (str): The path to the Excel file.

        Returns:
            List[Document]: A list of Document objects containing the extracted text in a single string.
        """
        data = pd.read_excel(path)
        structured_text = data.to_csv(index=False)
        return [Document(page_content=structured_text, metadata={ "source" : path })]
    
    def chunk_docs(self, raw_docs: list) -> list:
        """
        Break down a list of raw documents into smaller chunks of text.

        Each document is split into chunks of text with a maximum size of 512 characters.
        Each chunk has a 128 character overlap with the previous one to account for
        context.

        Args:
            raw_docs (List[Document]): The list of raw Document objects to be chunked.

        Returns:
            List[Document]: A list of Document objects containing the chunked text.
        """
        text_processor = RecursiveCharacterTextSplitter(
            chunk_size=512,
            chunk_overlap=128,
            add_start_index=True
        )
        return text_processor.split_documents(raw_docs)
    
    def index_docs(self, doc_chunk: list):
        """
        Index a list of Document objects into the vector database.

        Args:
            doc_chunk (List[Document]): The list of Document objects to be indexed.
        """
        self.document_vector_db.add_documents(doc_chunk)

    def find_related_docs(self, query: str) -> list:
        """
        Find the top 5 related documents to the given query.

        The method uses the vector database to search for the most similar documents
        to the query string. The result is a list of Document objects that are most
        semantically similar to the query.

        Args:
            query (str): The query string to search for related documents.

        Returns:
            List[Document]: A list of the top 5 most related documents.
        """
        return self.document_vector_db.similarity_search(query, k=5)
    
    def invoke(self, user_query: str, ctx_docs: list) -> str:
        """
        Use the given context documents to answer the given user query.

        The method combines the context documents into a single string and
        uses the vector database to search for the most similar documents
        to the query string. The result is a string that is the answer to the
        query based on the context.

        Args:
            user_query (str): The query string to answer.
            ctx_docs (list): A list of Document objects that provide context for the query.

        Returns:
            str: The answer to the query based on the context.
        """
        ctx_text = "\n\n".join([doc.page_content for doc in ctx_docs])
        conversation_prompt = ChatPromptTemplate.from_template(self.PROMPT_TEMPLATE)
        response_chain = conversation_prompt | self.model
        return response_chain.invoke({ "user_query": user_query, "document_context": ctx_text })